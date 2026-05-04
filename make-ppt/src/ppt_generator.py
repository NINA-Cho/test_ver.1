"""PPT 생성 엔진 - python-pptx로 .pptx 파일 생성."""
import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

DEFAULT_SPEC = {
    "background_color": "#FFFFFF",
    "accent_color": "#0066CC",
    "secondary_color": "#F5F5F5",
    "title_color": "#1A1A1A",
    "body_color": "#444444",
    "layout_type": "left-aligned",
    "title_position": "top",
    "font_ratio": "2.5:1",
    "style": "corporate",
}


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_rect(slide, left, top, width, height, color: str):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    return shape


def _add_text(
    slide, text: str, left, top, width, height,
    font_size: int, bold: bool, color: str,
    align=PP_ALIGN.LEFT, wrap: bool = True,
) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    p.alignment = align


def _add_bullets(
    slide, bullets: list, left, top, width, height,
    font_size: int, color: str,
) -> None:
    if not bullets:
        return
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = _rgb(color)
        p.space_after = Pt(8)


def _render_title(slide, info: dict, s: dict) -> None:
    _set_bg(slide, s["background_color"])
    _add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.08), s["accent_color"])
    _add_text(
        slide, info.get("title", ""),
        Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.6),
        font_size=42, bold=True, color=s["title_color"],
    )
    bullets = info.get("content", [])
    if bullets:
        _add_text(
            slide, bullets[0],
            Inches(0.8), Inches(4.1), Inches(9), Inches(1.0),
            font_size=20, bold=False, color=s["body_color"],
        )


def _render_section(slide, info: dict, s: dict) -> None:
    bg = s["secondary_color"] if s["secondary_color"] != s["background_color"] else "#F0F4F8"
    _set_bg(slide, bg)
    _add_rect(slide, Inches(0), Inches(3.55), SLIDE_W, Inches(0.06), s["accent_color"])
    _add_text(
        slide, info.get("title", ""),
        Inches(1), Inches(2.8), Inches(11), Inches(1.2),
        font_size=36, bold=True, color=s["accent_color"], align=PP_ALIGN.CENTER,
    )


def _render_content(slide, info: dict, s: dict) -> None:
    _set_bg(slide, s["background_color"])
    layout = s.get("layout_type", "left-aligned")
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), s["accent_color"])

    if layout == "two-column":
        _add_text(
            slide, info.get("title", ""),
            Inches(0.5), Inches(0.2), Inches(12.3), Inches(1.0),
            font_size=26, bold=True, color=s["title_color"],
        )
        bullets = info.get("content", [])
        mid = max(1, len(bullets) // 2)
        _add_bullets(slide, bullets[:mid], Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5), 17, s["body_color"])
        _add_bullets(slide, bullets[mid:], Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.5), 17, s["body_color"])
    else:
        title_align = PP_ALIGN.CENTER if layout == "centered" else PP_ALIGN.LEFT
        _add_text(
            slide, info.get("title", ""),
            Inches(0.6), Inches(0.2), Inches(12.1), Inches(1.0),
            font_size=28, bold=True, color=s["title_color"], align=title_align,
        )
        _add_bullets(
            slide, info.get("content", []),
            Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
            font_size=19, color=s["body_color"],
        )


def _render_closing(slide, info: dict, s: dict) -> None:
    _set_bg(slide, s["background_color"])
    _add_rect(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.08), s["accent_color"])
    _add_text(
        slide, info.get("title", "감사합니다"),
        Inches(1), Inches(2.8), Inches(11), Inches(1.5),
        font_size=44, bold=True, color=s["title_color"], align=PP_ALIGN.CENTER,
    )
    bullets = info.get("content", [])
    if bullets:
        _add_text(
            slide, bullets[0],
            Inches(1), Inches(4.5), Inches(11), Inches(1.0),
            font_size=18, bold=False, color=s["body_color"], align=PP_ALIGN.CENTER,
        )


RENDERERS = {
    "title": _render_title,
    "section": _render_section,
    "content": _render_content,
    "closing": _render_closing,
}


def create_ppt(
    slide_data: dict,
    design_spec: dict = None,
    output_path: str = None,
) -> str:
    """slide_data와 design_spec을 받아 .pptx 파일을 생성하고 경로를 반환."""
    spec = {**DEFAULT_SPEC, **(design_spec or {})}

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    for info in slide_data.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        renderer = RENDERERS.get(info.get("type", "content"), _render_content)
        renderer(slide, info, spec)
        notes = info.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    if not output_path:
        raw_title = slide_data.get("title", "presentation")
        safe = re.sub(r"[^\w\s-]", "", raw_title).strip().replace(" ", "_")
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = f"/tmp/{safe}_{ts}.pptx"

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
